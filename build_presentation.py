import sys
import copy
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

BASE = r'C:\Users\viach\Desktop\my-shop'
SRC = f'{BASE}\\reports\\presentation_gamma.pptx'
OUT = f'{BASE}\\reports\\Презентація_Раделицький_КНМС41.pptx'

SCREENSHOTS = f'{BASE}\\screenshots'
DIAGRAMS = f'{BASE}\\.claude\\worktrees\\stoic-antonelli-a1c284\\diagrams_out'

prs = Presentation(SRC)

def remove_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)

def replace_image(slide, old_shape, img_path):
    left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
    remove_shape(slide, old_shape)
    slide.shapes.add_picture(img_path, left, top, width, height)

def find_shapes_by_name(slide, names):
    return {s.name: s for s in slide.shapes if s.name in names}

# ============================================================
# Slide 6 (index 5): Use-case diagram placeholder
# ============================================================
slide6 = prs.slides[5]
shapes6 = {s.name: s for s in slide6.shapes}

# Replace placeholder image with customer use-case diagram
if 'Image 0' in shapes6:
    replace_image(slide6, shapes6['Image 0'], f'{DIAGRAMS}\\usecase_customer.png')

# ============================================================
# Slide 7 (index 6): DB structure — placeholder image
# We don't have a dedicated ER diagram, so leave as is or
# replace with admin use-case (it fits the slide about structure)
# Actually, let's leave the Gamma-generated placeholder for slide 7
# ============================================================

# ============================================================
# Slide 8 (index 7): Client interface screenshots
# Replace 4 large placeholders with actual screenshots
# Remove "Місце для скріншоту" boxes and small icons
# ============================================================
slide8 = prs.slides[7]

# Collect shapes by name
s8 = {s.name: s for s in slide8.shapes}

# Map: large placeholder image name -> screenshot file
slide8_map = {
    'Image 0': f'{SCREENSHOTS}\\ui_01_homepage.png',      # Головна
    'Image 2': f'{SCREENSHOTS}\\ui_02_catalog.png',        # Каталог
    'Image 4': f'{SCREENSHOTS}\\ui_03_product.png',        # Товар
    'Image 6': f'{SCREENSHOTS}\\ui_06_cart.png',           # Кошик
}

# Replace large placeholder images
for img_name, img_path in slide8_map.items():
    if img_name in s8:
        replace_image(slide8, s8[img_name], img_path)

# Remove small icons, "Місце для скріншоту" text boxes and their bg shapes
to_remove_8 = ['Image 1', 'Image 3', 'Image 5', 'Image 7',  # small icons
               'Text 5', 'Text 9', 'Text 13', 'Text 17',     # "Місце для скріншоту" text
               'Shape 4', 'Shape 8', 'Shape 12', 'Shape 16']  # background shapes
# Refresh shapes list after replacements
for s in list(slide8.shapes):
    if s.name in to_remove_8:
        remove_shape(slide8, s)

# Also remove the "студент вставляє вручну" subtitle
for s in list(slide8.shapes):
    if s.name == 'Text 1' and hasattr(s, 'text_frame'):
        if 'студент' in s.text_frame.text.lower():
            remove_shape(slide8, s)

# ============================================================
# Slide 9 (index 8): Checkout + Admin screenshots
# Replace 5 large placeholders with actual screenshots
# ============================================================
slide9 = prs.slides[8]
s9 = {s.name: s for s in slide9.shapes}

slide9_map = {
    'Image 0': f'{SCREENSHOTS}\\ui_07_checkout.png',           # Checkout
    'Image 2': f'{SCREENSHOTS}\\ui_08_admin_dashboard.png',    # Dashboard
    'Image 4': f'{SCREENSHOTS}\\ui_09_admin_orders.png',       # Замовлення
    'Image 6': f'{SCREENSHOTS}\\ui_11_admin_promo.png',        # Промокоди
    'Image 8': f'{SCREENSHOTS}\\ui_10_admin_users.png',        # Користувачі
}

for img_name, img_path in slide9_map.items():
    if img_name in s9:
        replace_image(slide9, s9[img_name], img_path)

# Remove small icons and "Місце для скріншоту" items
to_remove_9 = ['Image 1', 'Image 3', 'Image 5', 'Image 7', 'Image 9',
               'Text 5', 'Text 9', 'Text 13', 'Text 17', 'Text 21',
               'Shape 4', 'Shape 8', 'Shape 12', 'Shape 16', 'Shape 20']
for s in list(slide9.shapes):
    if s.name in to_remove_9:
        remove_shape(slide9, s)

# Remove "студент вставляє вручну" subtitle on slide 9
for s in list(slide9.shapes):
    if s.name == 'Text 1' and hasattr(s, 'text_frame'):
        if 'студент' in s.text_frame.text.lower():
            remove_shape(slide9, s)

# ============================================================
# Save
# ============================================================
prs.save(OUT)
print(f'Saved: {OUT}')
