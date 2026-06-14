import sys
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = r'C:\Users\viach\Desktop\my-shop'
SRC = f'{BASE}\\reports\\presentation_gamma.pptx'
OUT = f'{BASE}\\reports\\Презентація_Раделицький_КНМС41.pptx'

SCREENSHOTS = f'{BASE}\\screenshots'
DIAGRAMS = f'{BASE}\\.claude\\worktrees\\stoic-antonelli-a1c284\\diagrams_out'

prs = Presentation(SRC)

def remove_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)

def remove_shapes_by_names(slide, names):
    for s in list(slide.shapes):
        if s.name in names:
            remove_shape(slide, s)

def remove_all_except(slide, keep_names):
    for s in list(slide.shapes):
        if s.name not in keep_names:
            remove_shape(slide, s)

def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

# ============================================================
# Slide 6 (index 5): Replace use-case diagram placeholder
# ============================================================
slide6 = prs.slides[5]
s = get_shape(slide6, 'Image 0')
if s:
    left, top, width, height = s.left, s.top, s.width, s.height
    remove_shape(slide6, s)
    slide6.shapes.add_picture(f'{DIAGRAMS}\\usecase_customer.png', left, top, width, height)

# ============================================================
# Slide 8 (index 7): Client screenshots — replace placeholders
# ============================================================
slide8 = prs.slides[7]

# Map placeholder images to screenshots
slide8_images = {
    'Image 0': f'{SCREENSHOTS}\\ui_01_homepage.png',
    'Image 2': f'{SCREENSHOTS}\\ui_02_catalog.png',
    'Image 4': f'{SCREENSHOTS}\\ui_03_product.png',
    'Image 6': f'{SCREENSHOTS}\\ui_06_cart.png',
}

for img_name, img_path in slide8_images.items():
    s = get_shape(slide8, img_name)
    if s:
        left, top, width, height = s.left, s.top, s.width, s.height
        remove_shape(slide8, s)
        slide8.shapes.add_picture(img_path, left, top, width, height)

# Remove "Місце для скріншоту" related shapes
remove_shapes_by_names(slide8, [
    'Image 1', 'Image 3', 'Image 5', 'Image 7',
    'Text 5', 'Text 9', 'Text 13', 'Text 17',
    'Shape 4', 'Shape 8', 'Shape 12', 'Shape 16',
])

# Remove "студент вставляє" text
for s in list(slide8.shapes):
    if hasattr(s, 'text_frame') and 'студент' in s.text_frame.text.lower():
        remove_shape(slide8, s)

# ============================================================
# Slide 9 (index 8): Admin screenshots
# Gamma made 5 narrow rows — rebuild with bigger images
# ============================================================
slide9 = prs.slides[8]

# Save title shape info before clearing
title_shape = get_shape(slide9, 'Text 0')
title_left = title_shape.left if title_shape else Emu(434132)
title_top = title_shape.top if title_shape else Emu(394395)

# Remove ALL shapes from slide 9 and rebuild
for s in list(slide9.shapes):
    remove_shape(slide9, s)

# Slide dimensions
SW = prs.slide_width   # 9144000 EMU = 10 inches
SH = prs.slide_height  # 5143500 EMU = 5.625 inches

margin_left = Emu(434132)
margin_right = Emu(434132)
content_width = SW - margin_left - margin_right

# Add title
from pptx.util import Emu
from pptx.dml.color import RGBColor

txBox = slide9.shapes.add_textbox(margin_left, Emu(300000), content_width, Emu(320000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Оформлення замовлення та адмін-панель"
run = p.runs[0]
run.font.size = Pt(22)
run.font.bold = False
run.font.color.rgb = RGBColor(0x1B, 0x1B, 0x27)
run.font.name = "Raleway"

# Layout: 2 columns x 3 rows, with labels
screenshots_9 = [
    ('Рис. 7. Оформлення замовлення', f'{SCREENSHOTS}\\ui_07_checkout.png'),
    ('Рис. 8. Dashboard адмін-панелі', f'{SCREENSHOTS}\\ui_08_admin_dashboard.png'),
    ('Рис. 9. Управління замовленнями', f'{SCREENSHOTS}\\ui_09_admin_orders.png'),
    ('Рис. 10. Управління промокодами', f'{SCREENSHOTS}\\ui_11_admin_promo.png'),
    ('Рис. 11. Управління користувачами', f'{SCREENSHOTS}\\ui_10_admin_users.png'),
]

# 2 columns layout
col_gap = Emu(120000)
img_width = (content_width - col_gap) // 2
img_height = Emu(int(img_width * 0.45))  # ~16:9 aspect
label_height = Emu(180000)
row_gap = Emu(50000)

start_y = Emu(700000)
col_x = [margin_left, margin_left + img_width + col_gap]

for idx, (label, img_path) in enumerate(screenshots_9):
    row = idx // 2
    col = idx % 2

    y = start_y + row * (img_height + label_height + row_gap)
    x = col_x[col]

    # Add image
    slide9.shapes.add_picture(img_path, x, y, img_width, img_height)

    # Add label below image
    lbl = slide9.shapes.add_textbox(x, y + img_height, img_width, label_height)
    tf = lbl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    run = p.runs[0]
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x3C, 0x39, 0x39)
    run.font.name = "Roboto"
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER

# ============================================================
# Save
# ============================================================
prs.save(OUT)
print(f'OK: {OUT}')
print(f'Slides: {len(prs.slides)}')
